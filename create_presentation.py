from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Create a Presentation object
    prs = Presentation()
    
    # Title slide
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_title.shapes.title
    subtitle = slide_title.placeholders[1]
    title.text = "Internship Review"
    subtitle.text = "A summary of my internship experience"

    # Content slide
    slide_content = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_content.shapes.title
    content = slide_content.placeholders[1]
    title.text = "My Internship Experience"
    content.text = "During my internship, I learned about ..."

    # Save the presentation
    prs.save('internship_review.pptx')

if __name__ == '__main__':
    create_presentation()